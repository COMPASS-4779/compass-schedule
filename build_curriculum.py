"""
PDF + PPTX → curriculum.json 変換スクリプト
- 問題/解答PDF: 第NN章/問題_h_eng_gra_NN_XXX『...』（演習|習熟|習得）.pdf
- 解説PPTX:    第NN章/成果物/第NN章_XXX_<項目名>_内容.pptx
"""
import os, re, json, sys, subprocess
from concurrent.futures import ThreadPoolExecutor
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "curriculum.json")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

CHAPTER_TITLES = {
    1:"時制", 2:"受動態", 3:"助動詞", 4:"不定詞", 5:"動名詞",
    6:"分詞", 7:"仮定法", 8:"比較", 9:"関係詞", 10:"接続詞",
    11:"自動詞と他動詞", 12:"名詞", 13:"代名詞", 14:"前置詞",
    15:"形容詞・副詞", 16:"文の構造", 17:"強調・倒置・否定など"
}
CHAPTER_COLORS = {
    1:"#4caf50", 2:"#1976d2", 3:"#ff9800", 4:"#9c27b0", 5:"#00897b",
    6:"#5d4037", 7:"#e91e63", 8:"#3f51b5", 9:"#009688", 10:"#795548",
    11:"#607d8b", 12:"#f57c00", 13:"#7b1fa2", 14:"#0097a7",
    15:"#388e3c", 16:"#5e35b1", 17:"#c62828"
}

def pdftotext(path, layout=True):
    """layoutモードで縦並び順を保持"""
    try:
        args = ["pdftotext"]
        if layout: args.append("-layout")
        args += [path, "-"]
        r = subprocess.run(args, capture_output=True, timeout=15)
        return r.stdout.decode("utf-8", errors="ignore")
    except Exception as e:
        print(f"  [WARN] pdftotext fail: {path} -- {e}", file=sys.stderr)
        return ""

HEADER_PATTERNS = [
    re.compile(r'英文法\s*通常学習編'),
    re.compile(r'^\s*得点\s*$'),
    re.compile(r'^\s*100\s*$'),
    re.compile(r'^\s*ｗ\s*$'),
    re.compile(r'―解答―|―問題―'),
    re.compile(r'高 ?校 ?英 ?語'),
    re.compile(r'1?[0０]分確認テスト'),
    re.compile(r'Tomomitsu\s*Kimura'),
    re.compile(r'^\s*\d+\s*[-－,]\s*\d+(\s+演習|習熟|習得)?'),
]
def remove_header_lines(text):
    lines = text.split("\n")
    keep = []
    for line in lines:
        s = line.strip()
        if not s and not keep: continue
        if any(p.search(line) for p in HEADER_PATTERNS): continue
        # 異常空白の整理（layoutモードの空白を圧縮）
        line = re.sub(r' {3,}', '   ', line)
        keep.append(line.rstrip())
    # 末尾の空行除去
    while keep and not keep[-1].strip(): keep.pop()
    txt = "\n".join(keep).strip()
    txt = re.sub(r'\n{3,}', '\n\n', txt)
    return txt

def parse_filename(fn):
    m = re.match(r'(問題|解答)_h_eng_gra_(\d{2})_(\d{3})『(.+?)』（(演習|習熟|習得)）\.pdf$', fn)
    if not m: return None
    return {"kind":m.group(1), "ch":int(m.group(2)), "num":m.group(3), "name":m.group(4), "stage":m.group(5)}

STAGE_KEY = {"演習":"enshu","習熟":"shujuku","習得":"shutoku"}
STAGE_LABEL = {"enshu":"演習","shujuku":"習熟","shutoku":"習得"}

def scan_curriculum():
    chapters = {}
    for entry in sorted(os.listdir(ROOT)):
        m = re.match(r'^第(\d{2})章$', entry)
        if not m: continue
        ch_num = int(m.group(1))
        ch_dir = os.path.join(ROOT, entry)
        if not os.path.isdir(ch_dir): continue
        ch = chapters.setdefault(ch_num, {
            "id": f"ch{ch_num:02d}", "num": ch_num,
            "title": CHAPTER_TITLES.get(ch_num, f"第{ch_num}章"),
            "color": CHAPTER_COLORS.get(ch_num, "#666"),
            "items": {}
        })
        for fn in sorted(os.listdir(ch_dir)):
            info = parse_filename(fn)
            if not info: continue
            item = ch["items"].setdefault(info["num"], {
                "id": f"ch{ch_num:02d}_{info['num']}", "num": info["num"],
                "name": info["name"], "stages": {"enshu":{},"shujuku":{},"shutoku":{}},
                "slides": []
            })
            stage = item["stages"][STAGE_KEY[info["stage"]]]
            stage["problem_pdf" if info["kind"]=="問題" else "answer_pdf"] = os.path.join(entry, fn)
        # 成果物/*_内容.pptx を探して項目に紐付け
        seibutsu_dir = os.path.join(ch_dir, "成果物")
        if os.path.isdir(seibutsu_dir):
            for fn in sorted(os.listdir(seibutsu_dir)):
                m2 = re.match(r'^第(\d{2})章_(\d{3})_(.+)_内容\.pptx$', fn)
                if not m2: continue
                inum = m2.group(2)
                if inum in ch["items"]:
                    ch["items"][inum]["pptx_path"] = os.path.join(entry, "成果物", fn)
    return chapters

def extract_slides(pptx_path):
    """PPTXからスライド配列を生成"""
    if not Presentation: return []
    try:
        p = Presentation(pptx_path)
    except Exception as e:
        print(f"  [WARN] pptx open: {pptx_path} -- {e}", file=sys.stderr)
        return []
    slides = []
    for i, s in enumerate(p.slides):
        title = ""
        lines = []
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            for j, para in enumerate(sh.text_frame.paragraphs):
                t = "".join(r.text for r in para.runs).strip()
                if not t: continue
                if not title and j == 0 and not lines:
                    title = t
                else:
                    lines.append(t)
        # Slide1は表紙（章名と項目名）として扱う場合あり
        slides.append({
            "title": title or f"スライド {i+1}",
            "content": "\n".join(lines)
        })
    return slides

def estimate_question_count(text):
    if not text: return 0
    cnt = 0
    cnt += len(re.findall(r'[①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮⑯⑰⑱⑲⑳]', text))
    cnt += len(re.findall(r'\([0-9]{1,2}\)', text))
    cnt += len(re.findall(r'[⑴⑵⑶⑷⑸⑹⑺⑻⑼⑽⑾⑿⒀⒁⒂⒃⒄⒅⒆⒇]', text))
    return cnt

def extract_one(args):
    """並列ワーカー用"""
    item, stage_key, fld = args
    s = item["stages"][stage_key]
    if fld in s:
        full = os.path.join(ROOT, s[fld])
        txt = pdftotext(full)
        out_field = "problem_text" if fld == "problem_pdf" else "answer_text"
        s[out_field] = remove_header_lines(txt)
    return None

def extract_all_parallel(chapters, max_workers=8):
    tasks = []
    for ch in chapters.values():
        for it in ch["items"].values():
            for sk in ("enshu","shujuku","shutoku"):
                for fld in ("problem_pdf","answer_pdf"):
                    if fld in it["stages"][sk]:
                        tasks.append((it, sk, fld))
    print(f"  PDF tasks: {len(tasks)}, workers: {max_workers}")
    done = 0
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        for _ in ex.map(extract_one, tasks):
            done += 1
            if done % 100 == 0:
                print(f"    {done}/{len(tasks)}", file=sys.stderr)
    # PPTX 抽出（直列で十分速い）
    p_done = 0
    for ch in chapters.values():
        for it in ch["items"].values():
            if "pptx_path" in it:
                full = os.path.join(ROOT, it["pptx_path"])
                it["slides"] = extract_slides(full)
                p_done += 1
    print(f"  PPTX extracted: {p_done}")

def finalize(chapters):
    out = {"meta":{"version":"3.0","passMark":85}, "chapters":[]}
    for n in sorted(chapters.keys()):
        ch = chapters[n]
        items_list = []
        for inum in sorted(ch["items"].keys()):
            it = ch["items"][inum]
            stages_out = {}
            for sk in ("enshu","shujuku","shutoku"):
                s = it["stages"][sk]
                stages_out[sk] = {
                    "label": STAGE_LABEL[sk],
                    "problem_text": s.get("problem_text",""),
                    "answer_text": s.get("answer_text",""),
                    # 元PDFへの相対パス（ブラウザで開いたとき ../第NN章/... で参照可能）
                    "problem_pdf": "../" + s.get("problem_pdf","").replace("\\","/") if s.get("problem_pdf") else "",
                    "answer_pdf":  "../" + s.get("answer_pdf","").replace("\\","/") if s.get("answer_pdf") else "",
                    "q_count": estimate_question_count(s.get("problem_text",""))
                }
            items_list.append({
                "id": it["id"], "num": it["num"], "name": it["name"],
                "slides": it.get("slides", []),
                "pptx_path": ("../"+it["pptx_path"].replace("\\","/")) if it.get("pptx_path") else "",
                "stages": stages_out
            })
        out["chapters"].append({"id":ch["id"],"num":ch["num"],"title":ch["title"],"color":ch["color"],"items":items_list})
    return out

def process_chapters(ch_nums):
    chapters = scan_curriculum()
    sub = {n: chapters[n] for n in ch_nums if n in chapters}
    print(f"== extract chapters {ch_nums} ==")
    extract_all_parallel(sub, max_workers=8)
    final = finalize(sub)
    for n in ch_nums:
        out_path = os.path.join(os.path.dirname(OUT), f"ch{n:02d}.json")
        sub_out = {"meta":final["meta"],"chapters":[c for c in final["chapters"] if c["num"]==n]}
        with open(out_path, "w", encoding="utf-8") as fp:
            json.dump(sub_out, fp, ensure_ascii=False, indent=1)
        print(f"  saved ch{n:02d}.json (items={len(sub_out['chapters'][0]['items']) if sub_out['chapters'] else 0})")

def merge_all():
    out = {"meta":{"version":"3.0","passMark":85}, "chapters":[]}
    for n in range(1,18):
        p = os.path.join(os.path.dirname(OUT), f"ch{n:02d}.json")
        if os.path.exists(p):
            with open(p, encoding="utf-8") as fp:
                d = json.load(fp)
            out["chapters"].extend(d["chapters"])
    with open(OUT, "w", encoding="utf-8") as fp:
        json.dump(out, fp, ensure_ascii=False, indent=1)
    sz = os.path.getsize(OUT)/1024
    print(f"merged \u2192 {OUT} (chapters: {len(out['chapters'])}, {sz:.1f}KB)")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        cmd = sys.argv[1]
        if cmd == "merge":
            merge_all()
        elif cmd == "all":
            process_chapters(list(range(1,18)))
            merge_all()
        else:
            process_chapters([int(a) for a in sys.argv[1:]])
    else:
        print("usage: build_curriculum.py [N N N | all | merge]")
